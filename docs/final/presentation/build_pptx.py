#!/usr/bin/env python3
"""Assemble slides.pptx from the beamer page images + the speaker notes in slides.md.

Each .pptx slide is one full-bleed image of the beamer (Singapore/dolphin) page, and
the speaker notes are parsed from the `::: notes` blocks in slides.md (each note begins
with the presenter's name). This is how we get the beamer look inside an editable-notes
.pptx — pandoc does NOT apply beamer themes to its native pptx output.

Usage:  python3 build_pptx.py <dir-with-slide-PNGs>
        (build.sh generates the PNGs with pdftoppm and calls this script)
"""
import glob
import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
MD = os.path.join(HERE, "slides.md")
OUT = os.path.join(HERE, "slides.pptx")
PNG_DIR = sys.argv[1] if len(sys.argv) > 1 else "/tmp/fl_pngs"

# --- parse slides.md: ordered frames (# and ##) with their ::: notes block ---
lines = open(MD, encoding="utf-8").read().split("\n")
i = 0
if lines and lines[0].strip() == "---":
    i = 1
    while i < len(lines) and lines[i].strip() != "---":
        i += 1
    i += 1

elements = []  # one per frame: {"note": str|None}
cur = None
in_notes = False
buf = []


def flush():
    global cur, buf
    if cur is not None:
        note = "\n".join(buf).strip().replace("**", "")  # drop markdown bold
        cur["note"] = note or None
        elements.append(cur)
    buf = []


for line in lines[i:]:
    s = line.strip()
    if line.startswith("## "):
        flush(); cur = {"note": None}; in_notes = False
    elif line.startswith("# "):
        flush(); cur = {"note": None}; in_notes = False
    elif s == "::: notes":
        in_notes = True
    elif s == ":::" and in_notes:
        in_notes = False
    elif in_notes:
        buf.append(line)
flush()

# page 1 = beamer auto title page (no note); then each frame in document order
notes_by_page = [None] + [el["note"] for el in elements]

pngs = sorted(glob.glob(os.path.join(PNG_DIR, "slide-*.png")))
if not pngs:
    sys.exit(f"ERROR: no PNGs found in {PNG_DIR} (run build.sh, which generates them)")
if len(notes_by_page) != len(pngs):
    print(f"WARNING: page/note mismatch ({len(notes_by_page)} notes vs {len(pngs)} pages) — notes may misalign")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

n_notes = 0
for idx, png in enumerate(pngs):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    note = notes_by_page[idx] if idx < len(notes_by_page) else None
    if note:
        slide.notes_slide.notes_text_frame.text = note
        n_notes += 1

prs.save(OUT)
print(f"OK  {os.path.relpath(OUT)}  |  {len(pngs)} slides, {n_notes} with speaker notes")
