#!/usr/bin/env python3
"""
build_beautiful.py — FoodLens beautiful PPTX builder
Reads slides.md and outputs a professionally designed PPTX.

Usage:
  python3 build_beautiful.py
  python3 build_beautiful.py --input slides.md --output slides_beautiful.pptx
"""
import re
import sys
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree
except ImportError:
    print("Install: pip install python-pptx lxml")
    sys.exit(1)

# ── PALETTE ────────────────────────────────────────────────────────────────────
CREAM    = RGBColor(0xF3, 0xEE, 0xD8)
OLIVE    = RGBColor(0x2C, 0x50, 0x18)
OLIVE_LT = RGBColor(0x4A, 0x7A, 0x2A)
AMBER    = RGBColor(0xC8, 0x90, 0x2A)
INK      = RGBColor(0x1B, 0x1A, 0x0E)
SURFACE  = RGBColor(0xEC, 0xE7, 0xD4)
MUTED    = RGBColor(0x6B, 0x65, 0x50)
WHITE    = RGBColor(0xFA, 0xF8, 0xF0)
BORDER   = RGBColor(0xC8, 0xC2, 0xA8)
RED_TAG  = RGBColor(0xC0, 0x39, 0x2B)

# ── BASE DIR (resolved in build(), used by image helpers) ─────────────────────
_BASE_DIR = Path('.')

# ── FONTS ──────────────────────────────────────────────────────────────────────
SERIF = "Georgia"
SANS  = "Calibri"
MONO  = "Courier New"

# ── SLIDE DIMENSIONS 16:9 ─────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
MX = Inches(0.75)   # horizontal margin
MY = Inches(0.6)    # vertical margin
CW = W - 2 * MX    # content width
CH = H - 2 * MY    # content height


# ── HELPER: set slide background ───────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── HELPER: add a rectangle shape ──────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


# ── HELPER: add text box ────────────────────────────────────────────────────────
def add_textbox(slide, x, y, w, h, text, font_name=SANS, size=18,
                bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT,
                wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox, tf


# ── HELPER: add multiline text frame ───────────────────────────────────────────
def add_content_box(slide, x, y, w, h):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return txBox, tf


def add_para(tf, text, font_name=SANS, size=16, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, space_before=0, space_after=0,
             indent=0):
    from pptx.oxml.ns import qn as _qn
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ── IMAGE HELPERS ──────────────────────────────────────────────────────────────
def parse_image(line):
    """Parse '![](path){width=X%}' -> (path, width_fraction) or None."""
    m = re.match(r'!\[.*?\]\(([^)]+)\)\s*(?:\{[^}]*width=(\d+)%[^}]*\})?', line.strip())
    if m:
        path = m.group(1).strip()
        width_pct = int(m.group(2)) / 100 if m.group(2) else 0.70
        return path, width_pct
    return None


def add_image_to_slide(slide, img_rel_path, width_pct, y_start):
    """Insert image centered on slide. Returns bottom y-coordinate (EMU)."""
    full_path = _BASE_DIR / img_rel_path
    if not full_path.exists():
        print(f"  [WARN] image not found: {full_path}")
        return y_start
    img_w = int(CW * width_pct)
    x = MX + (CW - img_w) // 2
    pic = slide.shapes.add_picture(str(full_path), x, y_start, width=img_w)
    return y_start + pic.height


# ── MARKDOWN PARSER ─────────────────────────────────────────────────────────────
def strip_md_bold(text):
    """Remove **bold** markers, return clean text + bold flag."""
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text)


def parse_inline(text):
    """Return list of (text, bold, italic, mono) tuples."""
    parts = []
    i = 0
    while i < len(text):
        # code
        m = re.match(r'`([^`]+)`', text[i:])
        if m:
            parts.append((m.group(1), False, False, True))
            i += len(m.group(0))
            continue
        # bold+italic
        m = re.match(r'\*\*\*(.+?)\*\*\*', text[i:])
        if m:
            parts.append((m.group(1), True, True, False))
            i += len(m.group(0))
            continue
        # bold
        m = re.match(r'\*\*(.+?)\*\*', text[i:])
        if m:
            parts.append((m.group(1), True, False, False))
            i += len(m.group(0))
            continue
        # italic
        m = re.match(r'\*(.+?)\*', text[i:])
        if m:
            parts.append((m.group(1), False, True, False))
            i += len(m.group(0))
            continue
        # normal char
        if parts and not parts[-1][1] and not parts[-1][2] and not parts[-1][3]:
            parts[-1] = (parts[-1][0] + text[i], False, False, False)
        else:
            parts.append((text[i], False, False, False))
        i += 1
    return parts


def add_rich_para(tf, text, size=16, color=INK, align=PP_ALIGN.LEFT,
                  space_before=0, space_after=4, default_bold=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    parts = parse_inline(text)
    for chunk, bold, italic, mono in parts:
        run = p.add_run()
        run.text = chunk
        run.font.name = MONO if mono else SANS
        run.font.size = Pt(size)
        run.font.bold = bold or default_bold
        run.font.italic = italic
        run.font.color.rgb = MUTED if mono else color
    return p


def parse_slides_md(content):
    """Parse slides.md into a list of slide dicts."""
    # Strip YAML frontmatter
    content = re.sub(r'^---\n.*?\n---\n', '', content, flags=re.DOTALL)

    # Extract all notes blocks before splitting
    notes_map = {}  # We'll attach notes after splitting

    slides_raw = []
    current = {'type': None, 'title': '', 'content': [], 'notes': '', 'raw': ''}

    lines = content.split('\n')
    i = 0
    slide_idx = -1

    while i < len(lines):
        line = lines[i]

        # Section break
        if re.match(r'^# \d+\.', line):
            if current['type']:
                slides_raw.append(current)
            title = re.sub(r'^# ', '', line).strip()
            current = {'type': 'section', 'title': title, 'content': [],
                       'notes': '', 'raw': ''}
            slide_idx += 1
            i += 1
            continue

        # Slide heading
        if re.match(r'^## ', line):
            if current['type']:
                slides_raw.append(current)
            title = re.sub(r'^## ', '', line).strip()
            # Determine type from position
            slide_type = 'title' if not slides_raw and not current['type'] else 'content'
            current = {'type': slide_type, 'title': title, 'content': [],
                       'notes': '', 'raw': ''}
            slide_idx += 1
            i += 1
            continue

        # Notes block
        if line.strip() == '::: notes':
            notes_lines = []
            i += 1
            while i < len(lines) and lines[i].strip() != ':::':
                notes_lines.append(lines[i])
                i += 1
            current['notes'] = '\n'.join(notes_lines).strip()
            i += 1
            continue

        # Columns block — mark slide as twocol
        if line.strip() == '::: columns':
            current['type'] = 'twocol'
            cols = [[], []]
            col_idx = -1
            i += 1
            while i < len(lines) and lines[i].strip() != ':::':
                if re.match(r'^:::: \{\.column', lines[i]):
                    col_idx += 1
                elif lines[i].strip() == '::::':
                    pass
                else:
                    if col_idx >= 0 and col_idx < 2:
                        cols[col_idx].append(lines[i])
                i += 1
            current['cols'] = ['\n'.join(c).strip() for c in cols]
            i += 1
            continue

        # Regular content line
        if current['type']:
            current['content'].append(line)
        i += 1

    if current['type']:
        slides_raw.append(current)

    return slides_raw


def parse_bullets(text):
    """Parse text into list of (level, text) tuples for bullet items."""
    bullets = []
    for line in text.split('\n'):
        m = re.match(r'^(\s*)[-*]\s+(.*)', line)
        if m:
            level = len(m.group(1)) // 2
            bullets.append((level, m.group(2).strip()))
    return bullets


def parse_table(content_lines):
    """Parse markdown table into (headers, rows)."""
    table_lines = [l for l in content_lines if l.strip().startswith('|')]
    if not table_lines:
        return None, None
    headers = [c.strip() for c in table_lines[0].strip('|').split('|')]
    rows = []
    for line in table_lines[2:]:
        if '|' in line:
            rows.append([c.strip() for c in line.strip('|').split('|')])
    return headers, rows


# ── SLIDE RENDERERS ────────────────────────────────────────────────────────────

def render_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, CREAM)

    # Left olive accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, fill=OLIVE)

    # Eyebrow
    add_textbox(slide, MX + Inches(0.1), MY, CW, Inches(0.4),
                "UIB · HCI 11755 · Final Project · June 2026",
                font_name=MONO, size=10, color=MUTED)

    # Logo: Food + Lens in italic
    txBox = slide.shapes.add_textbox(MX + Inches(0.1), MY + Inches(0.5), CW, Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Food"
    r1.font.name = SERIF
    r1.font.size = Pt(72)
    r1.font.bold = False
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = "Lens"
    r2.font.name = SERIF
    r2.font.size = Pt(72)
    r2.font.italic = True
    r2.font.color.rgb = OLIVE

    # Tagline
    add_textbox(slide, MX + Inches(0.1), MY + Inches(2.5), Inches(8), Inches(0.9),
                "Transparent multi-objective food recommendations via contrastive XAI",
                font_name=SANS, size=18, color=MUTED)

    add_textbox(slide, MX + Inches(0.1), MY + Inches(3.15), Inches(8), Inches(0.5),
                "Health + Eco — shown together, never averaged, always explained.",
                font_name=SANS, size=15, italic=True, color=MUTED)

    # Separator line
    add_rect(slide, MX + Inches(0.1), MY + Inches(3.8), Inches(7), Inches(0.02), fill=BORDER)

    # Team names
    txBox2 = slide.shapes.add_textbox(MX + Inches(0.1), MY + Inches(4.0), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    names = ["Alejandro Bordón", "Soufyane Youbi", "Alejandro Rodríguez", "Pau Girón"]
    for j, name in enumerate(names):
        p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        r = p2.add_run()
        r.text = f"· {name}"
        r.font.name = SANS
        r.font.size = Pt(15)
        r.font.color.rgb = INK
        p2.space_after = Pt(2)

    # Course label bottom right
    add_textbox(slide, W - Inches(4), H - MY - Inches(0.5), Inches(3.2), Inches(0.5),
                "Course 11755 — Human-Computer Interaction\nMUSI-IA, UIB · June 2026",
                font_name=MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _add_notes(slide, data.get('notes', ''))


def render_section_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, OLIVE)

    title = data['title']
    # Extract number if present e.g. "1. Design process & rationale"
    m = re.match(r'^(\d+)\.\s+(.*)', title)
    num = m.group(1) if m else ''
    label = m.group(2) if m else title

    # Ghost large number
    if num:
        add_textbox(slide, W - Inches(3.5), H - Inches(3.5), Inches(3), Inches(3),
                    num, font_name=SERIF, size=180, color=WHITE, align=PP_ALIGN.RIGHT)
        # Make it transparent-ish by using a very light shade
        # (pptx doesn't support opacity in text easily, so we use a lighter color)
        # Actually let's use a slightly transparent-looking olive-light color
        txBox = slide.shapes[-1]
        tf = txBox.text_frame
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x4A, 0x7A, 0x2A)

    # Section label
    add_textbox(slide, MX, MY + Inches(0.3), CW, Inches(0.4),
                f"Section {num}" if num else "Section",
                font_name=MONO, size=11, color=RGBColor(0xAA, 0xCC, 0x88))

    # Section title
    txBox = slide.shapes.add_textbox(MX, MY + Inches(0.9), Inches(9), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    parts = label.split(' & ')
    if len(parts) == 2:
        r1 = p.add_run()
        r1.text = parts[0] + '\n'
        r1.font.name = SERIF
        r1.font.size = Pt(48)
        r1.font.color.rgb = WHITE
        r1.font.bold = False
        r2 = p.add_run()
        r2.text = '& ' + parts[1]
        r2.font.name = SERIF
        r2.font.size = Pt(48)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(0xCC, 0xEE, 0xAA)
    else:
        r = p.add_run()
        r.text = label
        r.font.name = SERIF
        r.font.size = Pt(48)
        r.font.color.rgb = WHITE

    _add_notes(slide, data.get('notes', ''))


def render_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    title = data['title']
    all_lines = [l for l in data.get('content', []) if l.strip()]

    # Top olive strip
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill=OLIVE)

    # Title
    txBox = slide.shapes.add_textbox(MX, MY, CW, Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    parts = parse_inline(title)
    for chunk, bold, italic, mono in parts:
        r = p.add_run()
        r.text = chunk
        r.font.name = SERIF
        r.font.size = Pt(30)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = OLIVE if (bold or italic) else INK

    y_pos = MY + Inches(1.1)

    # Separate image lines from text lines
    img_lines  = [l for l in all_lines if parse_image(l)]
    text_lines = [l for l in all_lines if not parse_image(l)]

    if img_lines:
        img_path, width_pct = parse_image(img_lines[0])
        y_pos = add_image_to_slide(slide, img_path, width_pct, y_pos)
        y_pos += Inches(0.15)

    content_text = '\n'.join(text_lines)
    if '|---|' in content_text or '|---|' in content_text.replace('| --- |', '|---|'):
        render_table_content(slide, text_lines, y_pos)
    else:
        render_bullets_content(slide, text_lines, y_pos)

    _add_notes(slide, data.get('notes', ''))


def render_bullets_content(slide, content_lines, y_start):
    txBox = slide.shapes.add_textbox(MX, y_start, CW, H - y_start - MY)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True

    for line in content_lines:
        if not line.strip():
            continue
        # Numbered list item
        m_num = re.match(r'^(\d+)\.\s+(.*)', line)
        # Bullet
        m_bullet = re.match(r'^[-*]\s+(.*)', line)
        # ### sub-heading
        m_h3 = re.match(r'^###\s+(.*)', line)
        # Regular paragraph
        if m_h3:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            r = p.add_run()
            r.text = strip_md_bold(m_h3.group(1))
            r.font.name = SERIF
            r.font.size = Pt(18)
            r.font.color.rgb = OLIVE
            p.space_before = Pt(8)
            first = False
        elif m_num:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(4)
            # Number in olive
            r1 = p.add_run()
            r1.text = f"{m_num.group(1)}.  "
            r1.font.name = MONO
            r1.font.size = Pt(15)
            r1.font.color.rgb = OLIVE
            r1.font.bold = True
            # Content
            for chunk, bold, italic, mono in parse_inline(m_num.group(2)):
                r = p.add_run()
                r.text = chunk
                r.font.name = MONO if mono else SANS
                r.font.size = Pt(16)
                r.font.bold = bold
                r.font.italic = italic
                r.font.color.rgb = MUTED if mono else INK
            first = False
        elif m_bullet:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(5)
            # Bullet dot
            r1 = p.add_run()
            r1.text = "·  "
            r1.font.name = SANS
            r1.font.size = Pt(18)
            r1.font.color.rgb = OLIVE
            r1.font.bold = True
            # Content
            for chunk, bold, italic, mono in parse_inline(m_bullet.group(1)):
                r = p.add_run()
                r.text = chunk
                r.font.name = MONO if mono else SANS
                r.font.size = Pt(16)
                r.font.bold = bold
                r.font.italic = italic
                r.font.color.rgb = MUTED if mono else INK
            first = False
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(6)
            for chunk, bold, italic, mono in parse_inline(line.strip()):
                r = p.add_run()
                r.text = chunk
                r.font.name = MONO if mono else SANS
                r.font.size = Pt(15)
                r.font.bold = bold
                r.font.italic = italic
                r.font.color.rgb = MUTED if mono else INK
            first = False


def render_table_content(slide, content_lines, y_start):
    from pptx.util import Inches as I2
    headers, rows = parse_table(content_lines)
    if not headers or not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows)

    tbl_w = CW
    tbl_h = Inches(0.42) * (n_rows + 1)
    col_w = tbl_w // n_cols

    table = slide.shapes.add_table(n_rows + 1, n_cols, MX, y_start, tbl_w, tbl_h).table

    col_widths = []
    # First col narrower if it looks like an ID col
    if n_cols == 2:
        col_widths = [Inches(1.4), tbl_w - Inches(1.4)]
    else:
        col_widths = [tbl_w // n_cols] * n_cols
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = OLIVE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = strip_md_bold(h)
        r.font.name = SANS
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE

    # Data rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            parts = parse_inline(strip_md_bold(cell_text))
            for chunk, bold, italic, mono in parts:
                r = p.add_run()
                r.text = chunk
                r.font.name = MONO if mono else SANS
                r.font.size = Pt(11) if j > 0 else Pt(12)
                r.font.bold = bold or (j == 0)
                r.font.color.rgb = OLIVE if j == 0 else INK


def render_twocol_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    title = data['title']
    cols = data.get('cols', ['', ''])

    # Top strip
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill=OLIVE)

    # Title
    txBox = slide.shapes.add_textbox(MX, MY, CW, Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    for chunk, bold, italic, mono in parse_inline(title):
        r = p.add_run()
        r.text = chunk
        r.font.name = SERIF
        r.font.size = Pt(28)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = OLIVE if (bold or italic) else INK

    # Vertical divider
    col_w = (CW - Inches(0.3)) / 2
    add_rect(slide, MX + col_w + Inches(0.15), MY + Inches(1.1),
             Inches(0.02), H - MY - Inches(1.3), fill=BORDER)

    y_content = MY + Inches(1.15)
    h_content = H - y_content - MY

    for ci, col_text in enumerate(cols[:2]):
        x = MX + ci * (col_w + Inches(0.3))
        col_lines = col_text.split('\n')

        txBox2 = slide.shapes.add_textbox(x, y_content, col_w, h_content)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        first = True

        for line in col_lines:
            if not line.strip():
                continue
            m_h = re.match(r'^\*\*(.+)\*\*$', line.strip())
            m_b = re.match(r'^[-*]\s+(.*)', line)
            m_sub = re.match(r'^  [-*]\s+(.*)', line)

            if m_h:
                p2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
                r = p2.add_run()
                r.text = m_h.group(1)
                r.font.name = SERIF
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = OLIVE
                p2.space_before = Pt(6)
                p2.space_after = Pt(4)
                first = False
            elif m_b:
                p2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
                p2.space_before = Pt(4)
                r1 = p2.add_run()
                r1.text = "·  "
                r1.font.name = SANS
                r1.font.size = Pt(16)
                r1.font.color.rgb = OLIVE
                r1.font.bold = True
                for chunk, bold, italic, mono in parse_inline(m_b.group(1)):
                    r = p2.add_run()
                    r.text = chunk
                    r.font.name = MONO if mono else SANS
                    r.font.size = Pt(15)
                    r.font.bold = bold
                    r.font.italic = italic
                    r.font.color.rgb = MUTED if mono else INK
                first = False
            else:
                p2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
                for chunk, bold, italic, mono in parse_inline(line.strip()):
                    r = p2.add_run()
                    r.text = chunk
                    r.font.name = MONO if mono else SANS
                    r.font.size = Pt(14)
                    r.font.bold = bold
                    r.font.italic = italic
                    r.font.color.rgb = MUTED if mono else INK
                first = False

    _add_notes(slide, data.get('notes', ''))


def render_outline_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill=OLIVE)

    add_textbox(slide, MX, MY, CW, Inches(0.8),
                "Overview", font_name=MONO, size=11, color=MUTED)

    add_textbox(slide, MX, MY + Inches(0.4), Inches(6), Inches(1.0),
                "Outline", font_name=SERIF, size=36, color=INK)

    items = [l for l in data.get('content', []) if re.match(r'^\d+\.', l.strip())]
    y = MY + Inches(1.5)
    item_h = (H - y - MY) / max(len(items), 1)

    for i, item in enumerate(items):
        m = re.match(r'^(\d+)\.\s+(.*)', item.strip())
        if not m:
            continue
        num = m.group(1)
        text = m.group(2)

        # Number box
        add_rect(slide, MX, y + i * item_h, Inches(0.5), item_h - Inches(0.06), fill=SURFACE)
        add_textbox(slide, MX + Inches(0.05), y + i * item_h + Inches(0.04),
                    Inches(0.4), item_h,
                    num, font_name=MONO, size=11, color=OLIVE, align=PP_ALIGN.CENTER)

        # Separator line
        add_rect(slide, MX + Inches(0.52), y + i * item_h + item_h / 2,
                 CW - Inches(0.52), Inches(0.01), fill=BORDER)

        add_textbox(slide, MX + Inches(0.65), y + i * item_h + Inches(0.04),
                    CW - Inches(0.65), item_h,
                    text, font_name=SANS, size=17, color=INK)

    _add_notes(slide, data.get('notes', ''))


def render_close_slide(prs, data):
    """Special treatment for the closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill=OLIVE)

    add_textbox(slide, MX, MY, CW, Inches(0.5),
                "Next Steps", font_name=MONO, size=11, color=MUTED)

    add_textbox(slide, MX, MY + Inches(0.4), Inches(8), Inches(1.0),
                "Future work & close", font_name=SERIF, size=30, color=INK)

    content_lines = [l for l in data.get('content', []) if l.strip()]

    # Bullets in upper half
    txBox = slide.shapes.add_textbox(MX, MY + Inches(1.5), CW, Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in content_lines:
        m = re.match(r'^[-*]\s+(.*)', line)
        if not m:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = "·  "
        r1.font.name = SANS
        r1.font.size = Pt(16)
        r1.font.color.rgb = OLIVE
        r1.font.bold = True
        for chunk, bold, italic, mono in parse_inline(m.group(1)):
            r = p.add_run()
            r.text = chunk
            r.font.name = MONO if mono else SANS
            r.font.size = Pt(15)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = INK
        first = False

    # Bottom close block
    add_rect(slide, MX, H - MY - Inches(1.8), CW, Inches(0.02), fill=BORDER)

    close_box = slide.shapes.add_textbox(MX, H - MY - Inches(1.7), Inches(8), Inches(0.9))
    tf2 = close_box.text_frame
    p2 = tf2.paragraphs[0]
    r_food = p2.add_run()
    r_food.text = "Food"
    r_food.font.name = SERIF
    r_food.font.size = Pt(36)
    r_food.font.color.rgb = INK
    r_lens = p2.add_run()
    r_lens.text = "Lens"
    r_lens.font.name = SERIF
    r_lens.font.size = Pt(36)
    r_lens.font.italic = True
    r_lens.font.color.rgb = OLIVE

    add_textbox(slide, MX, H - MY - Inches(0.9), Inches(7), Inches(0.4),
                "Dual-axis · Contrastive · Honest   —   github.com/AlejandroB10/foodlens",
                font_name=MONO, size=10, color=MUTED)

    add_textbox(slide, MX, H - MY - Inches(0.55), Inches(7), Inches(0.5),
                "Thank you — questions?",
                font_name=SERIF, size=24, color=INK)

    _add_notes(slide, data.get('notes', ''))


def _add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    # Strip bold markers from speaker name
    clean = re.sub(r'\*\*([^*]+)\*\*:', r'\1:', notes_text)
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = clean


# ── MAIN BUILD ─────────────────────────────────────────────────────────────────

def build(input_path: str, output_path: str):
    global _BASE_DIR
    _BASE_DIR = Path(input_path).parent
    text = Path(input_path).read_text(encoding='utf-8')
    slides_data = parse_slides_md(text)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Ensure at least one blank layout
    from pptx.util import Inches as _I
    for i in range(7 - len(prs.slide_layouts)):
        pass  # layouts are already there

    for i, slide in enumerate(slides_data):
        stype = slide['type']
        title = slide['title'].lower()

        print(f"  Slide {i+1:02d}: [{stype:8s}] {slide['title'][:50]}")

        if stype == 'section':
            render_section_slide(prs, slide)
        elif stype == 'twocol':
            render_twocol_slide(prs, slide)
        elif i == 0 or stype == 'title':
            render_title_slide(prs, slide)
        elif 'outline' in title or title.strip().lower() == 'outline':
            render_outline_slide(prs, slide)
        elif 'future work' in title or 'close' in title:
            render_close_slide(prs, slide)
        else:
            # Check for table in content
            content_text = '\n'.join(slide.get('content', []))
            if re.search(r'\|.*\|.*\|', content_text) and '---' in content_text:
                slide['_has_table'] = True
            render_content_slide(prs, slide)

    prs.save(output_path)
    print(f"\n✓ Saved → {output_path}  ({len(slides_data)} slides)")


def main():
    here = Path(__file__).parent
    ap = argparse.ArgumentParser(description="FoodLens beautiful PPTX builder")
    ap.add_argument('--input',  default=str(here / 'slides.md'),              help='Input markdown file')
    ap.add_argument('--output', default=str(here / 'slides_beautiful.pptx'),  help='Output PPTX file')
    args = ap.parse_args()

    if not Path(args.input).exists():
        print(f"Error: {args.input} not found")
        sys.exit(1)

    print(f"Building {args.output} from {args.input}…")
    build(args.input, args.output)


if __name__ == '__main__':
    main()